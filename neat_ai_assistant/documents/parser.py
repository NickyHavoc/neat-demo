from typing import Sequence
import PyPDF2
import docx
import pandas as pd

from pathlib import Path
from pptx import Presentation


# ToDos:
# - add chunking here (otherwise long pages will throw context length errors)


class Parser:
    def __init__(self):
        pass

    def parse_pdf(self, file_path: Path) -> Sequence[str]:
        content = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                content.append(page.extract_text())
        return content

    def parse_xlsx(self, file_path: Path) -> Sequence[str]:
        content = []
        df = pd.read_excel(file_path)
        for _, row in df.iterrows():
            row_content = ' '.join(row.astype(str).values.tolist())
            content.append(row_content)
        return content

    def parse_word(self, file_path: Path) -> Sequence[str]:
        content = []
        doc = docx.Document(file_path)
        for paragraph in doc.paragraphs:
            content.append(paragraph.text)
        return content

    def parse_ppt(self, file_path: Path) -> Sequence[str]:
        content = []
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
        return content

    def parse_txt(self, file_path: Path) -> Sequence[str]:
        with open(file_path, 'r') as file:
            content = file.readlines()
        return content

    def parse_file(self, file_path: Path) -> Sequence[str]:
        file_extension = file_path.suffix.lower()
        if file_extension == '.pdf':
            content = self.parse_pdf(file_path)
        elif file_extension == '.xlsx':
            content = self.parse_xlsx(file_path)
        elif file_extension == '.docx':
            content = self.parse_word(file_path)
        elif file_extension == '.pptx':
            content = self.parse_ppt(file_path)
        elif file_extension == '.txt':
            content = self.parse_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_path.suffix}")
        return content
